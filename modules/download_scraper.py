import contextlib
import csv
import os
import re
import shutil
import traceback
from collections.abc import Callable

import fitz  # PyMuPDF (já está no requirements.txt)
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyQt6.QtCore import QThread, pyqtSignal

from modules.config import (
    PALAVRAS_CHAVE,
    PERFIL_EDGE_DOWNLOAD,
    get_coupa_base_url,
    resolve_tesseract_executable,
)
from modules.playwright_pool import PlaywrightContextManager

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None  # type: ignore[assignment]
    Image = None  # type: ignore[assignment]

# O anexo do item do carrinho usa o MESMO componente visual (span.underline
# com aria-label "Anexo de arquivo de <nome>", dentro de um span.attachment-
# file) que a seção de Anexos (Justificativa) - não dá pra diferenciar pelo
# elemento em si. A diferença real, confirmada inspecionando o HTML real do
# Coupa, está no ancestral: o anexo do item do carrinho fica dentro de um
# container .lineAttachments, que não existe para os anexos de cima.
_SELETOR_ANEXOS = "span[aria-label^='Anexo de arquivo']"


class DownloadScraper:
    def __init__(self, requisicoes: list[str], pasta_download: str):
        # Nunca processa a mesma requisição duas vezes (ex: uma requisição com
        # 2 pedidos aparece 2x na lista importada da Aba 1) - baixar a mesma
        # página/anexo de novo é desperdício e ainda gera arquivo duplicado
        # (analisar_arquivo salva o 2º como "REQ_1.pdf").
        self.requisicoes = list(dict.fromkeys(r.strip() for r in requisicoes if r.strip()))
        self.pasta_download = pasta_download
        self.arquivos_salvos_na_execucao: list[str] = []
        self.requisicoes_sem_arquivos: list[str] = []
        self.cancelado = False
        self._log_callback: Callable[[str], None] | None = None

    @staticmethod
    def normalizar_texto(texto: str) -> str:
        texto = os.path.normpath(texto or "")
        texto = texto.replace("\\", " ")
        texto = texto.replace("/", " ")
        texto = texto.strip()
        return re.sub(r"\s+", " ", texto).strip().lower()

    def contem_de_acordo(self, texto: str) -> bool:
        texto = texto.replace("_", " ").replace("-", " ")
        return "de acordo" in self.normalizar_texto(texto)

    def _extrair_texto_log(self, msg: str) -> None:
        """Melhoria 8: loga erro via _log_callback ou logger padrão."""
        try:
            if self._log_callback:
                self._log_callback(msg)
            else:
                import logging
                logging.getLogger(__name__).warning(msg)
        except Exception:
            pass  # best-effort: nem o log de erro deve derrubar o fluxo

    def _ocr_pdf(self, caminho_arquivo: str) -> str:
        """OCR de PDF escaneado (sem camada de texto embutida), via Tesseract.

        Usado como fallback quando a extração normal (fitz) não encontra
        texto nenhum - típico de PDF que é só uma imagem escaneada. No
        executável instalado, usa o Tesseract empacotado junto (ver
        resolve_tesseract_executable) - o usuário não precisa instalar nada.
        Em modo desenvolvimento, cai no Tesseract do PATH do sistema, se
        houver. Se não estiver disponível de nenhuma forma, loga um aviso e
        retorna "" sem quebrar o fluxo - o arquivo cai na validação por nome,
        como antes do OCR existir.
        """
        if pytesseract is None or Image is None:
            self._extrair_texto_log(
                "[DownloadScraper] OCR indisponível: pacote pytesseract/Pillow não instalado."
            )
            return ""

        tesseract_empacotado = resolve_tesseract_executable()
        if tesseract_empacotado:
            pytesseract.pytesseract.tesseract_cmd = tesseract_empacotado

        texto_paginas = []
        try:
            with fitz.open(caminho_arquivo) as pdf:
                for pagina in pdf:
                    pixmap = pagina.get_pixmap(dpi=300)
                    imagem = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
                    texto_paginas.append(pytesseract.image_to_string(imagem, lang="por"))
        except pytesseract.TesseractNotFoundError:
            self._extrair_texto_log(
                "[DownloadScraper] OCR indisponível: Tesseract não encontrado nesta máquina. "
                "Instale em https://github.com/UB-Mannheim/tesseract/wiki (inclua o pacote de "
                "idioma Português) e adicione ao PATH, ou configure "
                "pytesseract.pytesseract.tesseract_cmd."
            )
            return ""
        except Exception as e:
            self._extrair_texto_log(
                f"[DownloadScraper] Erro no OCR de '{os.path.basename(caminho_arquivo)}': {str(e)}"
            )
            return ""
        return "\n".join(texto_paginas)

    def extrair_texto(self, caminho_arquivo: str) -> str:
        """Extrai texto de arquivos PDF, DOCX, TXT, CSV, XLSX, PPTX.

        Melhoria 8: logging via callback em vez de print().
        """
        extensao = os.path.splitext(caminho_arquivo)[1].lower()
        texto = ""
        try:
            if extensao == ".pdf":
                with fitz.open(caminho_arquivo) as pdf:
                    for pagina in pdf:
                        texto += pagina.get_text() or ""
                if not texto.strip():
                    # PDF sem texto embutido (provavelmente escaneado/imagem) - tenta OCR.
                    texto = self._ocr_pdf(caminho_arquivo)
            elif extensao == ".docx":
                documento = Document(caminho_arquivo)
                for paragrafo in documento.paragraphs:
                    texto += paragrafo.text + "\n"
            elif extensao == ".txt":
                with open(caminho_arquivo, encoding="utf-8", errors="ignore") as f:
                    texto = f.read()
            elif extensao == ".csv":
                with open(caminho_arquivo, newline="", encoding="utf-8", errors="ignore") as f:
                    leitor = csv.reader(f)
                    for linha in leitor:
                        texto += " ".join(map(str, linha)) + "\n"
            elif extensao == ".xlsx":
                workbook = load_workbook(caminho_arquivo, data_only=True)
                for aba in workbook.worksheets:
                    for linha in aba.iter_rows(values_only=True):
                        for celula in linha:
                            if celula is not None:
                                texto += str(celula) + " "
            elif extensao == ".pptx":
                apresentacao = Presentation(caminho_arquivo)
                for slide in apresentacao.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texto += shape.text + "\n"
        except Exception as e:
            self._extrair_texto_log(
                f"[DownloadScraper] Erro ao extrair texto de "
                f"'{os.path.basename(caminho_arquivo)}': {str(e)}"
            )
            return ""
        return texto.lower()

    _PADRAO_VALOR_TOTAL = re.compile(r"valor total\D{0,12}?(\d{1,3}(?:\.\d{3})*,\d{2})")

    def extrair_valor_total(self, texto_normalizado: str) -> float | None:
        """Extrai o valor de "Valor total: R$ X,XX" do texto (já normalizado -
        minúsculo, ver normalizar_texto), convertendo pro formato numérico.

        Usado para desempatar quando uma requisição tem mais de um documento
        válido como orçamento (ex: uma versão com o valor antigo e outra já
        com o valor negociado) - ver _escolher_melhor_candidato. Retorna None
        se não encontrar o padrão (o arquivo ainda conta como candidato, só
        não entra na comparação por valor).
        """
        match = self._PADRAO_VALOR_TOTAL.search(texto_normalizado)
        if not match:
            return None
        valor_str = match.group(1).replace(".", "").replace(",", ".")
        try:
            return float(valor_str)
        except ValueError:
            return None

    def _prioridade_palavra_chave(self, texto_normalizado: str) -> int:
        """Índice (em PALAVRAS_CHAVE) da palavra-chave de MAIOR prioridade
        encontrada no texto - quanto menor o índice, maior a prioridade (ver
        ordem/comentário em PALAVRAS_CHAVE, modules/config.py). Usado por
        _manter_apenas_melhor_candidato como critério principal de desempate.

        Retorna len(PALAVRAS_CHAVE) (pior prioridade possível) se nada bater -
        não deveria acontecer com um candidato que já passou por
        analisar_arquivo, exceto no caso raro de um arquivo sem texto
        embutido que só validou pelo NOME original (essa informação já se
        perdeu quando chegamos aqui, só o conteúdo salvo é reconsultado).
        """
        for indice, palavra in enumerate(PALAVRAS_CHAVE):
            if self.normalizar_texto(palavra) in texto_normalizado:
                return indice
        return len(PALAVRAS_CHAVE)

    def arquivo_tem_formato_invalido(self, texto: str) -> bool:
        texto_normalizado = self.normalizar_texto(texto)
        if "de acordo exceto" in texto_normalizado:
            return True

        header_terms = ["de:", "para:", "assunto:", "cc:", "interno"]
        if all(term in texto_normalizado for term in header_terms):
            return True

        return "escopos ajustados" in texto_normalizado and "de acordo" in texto_normalizado

    def analisar_arquivo(
        self,
        caminho_arquivo: str,
        req_num: str,
        nome_original: str,
        log_callback=None,
    ) -> tuple[bool, str | None]:
        extensao = os.path.splitext(caminho_arquivo)[1]
        destino = os.path.normpath(os.path.join(self.pasta_download, f"{req_num}{extensao}"))

        contador = 1
        while os.path.exists(destino):
            destino = os.path.normpath(os.path.join(self.pasta_download, f"{req_num}_{contador}{extensao}"))
            contador += 1

        if self.contem_de_acordo(nome_original):
            return False, "nome"

        texto = self.extrair_texto(caminho_arquivo)
        texto_normalizado = self.normalizar_texto(texto)

        if texto_normalizado == "":
            nome_normalizado = self.normalizar_texto(nome_original)
            if any(self.normalizar_texto(palavra) in nome_normalizado for palavra in PALAVRAS_CHAVE):
                shutil.copy(caminho_arquivo, destino)
                self.arquivos_salvos_na_execucao.append(destino)
                return True, destino
            return False, "vazio"

        if self.arquivo_tem_formato_invalido(texto_normalizado):
            return False, "formato"

        for palavra in PALAVRAS_CHAVE:
            if self.normalizar_texto(palavra) in texto_normalizado:
                shutil.copy(caminho_arquivo, destino)
                self.arquivos_salvos_na_execucao.append(destino)
                return True, destino

        if any(self.normalizar_texto(palavra) in self.normalizar_texto(nome_original) for palavra in PALAVRAS_CHAVE):
            shutil.copy(caminho_arquivo, destino)
            self.arquivos_salvos_na_execucao.append(destino)
            return True, destino

        return False, "palavra"

    async def run(self, log_callback, progress_req_callback, progress_down_callback) -> bool:
        log_callback("Iniciando Edge com Perfil Isolado...")
        user_data_dir = str(PERFIL_EDGE_DOWNLOAD)  # Item 16: path centralizado em config.py
        PERFIL_EDGE_DOWNLOAD.mkdir(parents=True, exist_ok=True)

        try:
            async with PlaywrightContextManager(user_data_dir=user_data_dir) as context:
                return await self._processar(context, log_callback, progress_req_callback, progress_down_callback)
        except Exception as e:
            log_callback(f"ERRO CRÍTICO AO INICIAR EDGE: {str(e)}")
            return False

    def _manter_apenas_melhor_candidato(self, caminhos_salvos: list[str], req: str) -> None:
        """Quando mais de um documento validou como orçamento pra mesma
        requisição, mantém só o melhor e apaga os demais.

        Critério PRINCIPAL: prioridade da palavra-chave encontrada (ver
        PALAVRAS_CHAVE em modules/config.py e _prioridade_palavra_chave) - um
        documento que bate em "orçamento" vence outro que só bate em "valor
        total", por exemplo, mesmo que este tenha preço menor.

        Critério de DESEMPATE (só entra quando 2+ candidatos batem na MESMA
        palavra-chave de maior prioridade): o de MENOR valor total (ver
        extrair_valor_total) - ex. uma versão com o valor antigo e outra já
        com o valor negociado, negociação normalmente reduz o preço, então o
        valor mais baixo é o vigente. Sem "valor total" extraível fica por
        último dentro do próprio grupo de prioridade (não descartado de
        cara, só perde a preferência frente a um valor real).
        """
        candidatos = []
        for caminho in caminhos_salvos:
            texto_normalizado = self.normalizar_texto(self.extrair_texto(caminho))
            prioridade = self._prioridade_palavra_chave(texto_normalizado)
            valor = self.extrair_valor_total(texto_normalizado)
            candidatos.append((prioridade, valor is None, valor if valor is not None else 0.0, caminho))
        candidatos.sort()

        prioridade_escolhida, sem_valor_escolhido, valor_escolhido, _ = candidatos[0]
        for _, _, _, caminho in candidatos[1:]:
            with contextlib.suppress(OSError):
                os.remove(caminho)
            self.arquivos_salvos_na_execucao.remove(caminho)

        palavra_escolhida = (
            PALAVRAS_CHAVE[prioridade_escolhida] if prioridade_escolhida < len(PALAVRAS_CHAVE) else None
        )
        houve_empate_de_prioridade = sum(1 for c in candidatos if c[0] == prioridade_escolhida) > 1

        if houve_empate_de_prioridade:
            contexto = f' entre os que bateram em "{palavra_escolhida}"' if palavra_escolhida else ""
            if sem_valor_escolhido:
                descricao_escolha = (
                    f"mantido o primeiro encontrado{contexto} "
                    "(não foi possível extrair o valor de nenhum)"
                )
            else:
                valor_str = f"R$ {valor_escolhido:,.2f}".replace(",", "_").replace(".", ",").replace("_", ".")
                descricao_escolha = f"mantido o de menor valor ({valor_str}){contexto}"
        elif palavra_escolhida:
            descricao_escolha = f'mantido o que bateu na palavra-chave de maior prioridade ("{palavra_escolhida}")'
        else:
            descricao_escolha = "mantido o primeiro encontrado"

        self._extrair_texto_log(
            f"ℹ️ #{req}: {len(caminhos_salvos)} documentos válidos encontrados - "
            f"{descricao_escolha}, os demais foram descartados."
        )

    async def _baixar_e_validar_anexos(
        self, page, elementos, req: str, progress_down_callback,
    ) -> int:
        """Baixa cada elemento de anexo e valida com analisar_arquivo (nome + conteúdo).

        Se mais de um anexo do grupo validar como orçamento, mantém só o
        melhor candidato (ver _manter_apenas_melhor_candidato) - o retorno já
        reflete isso (no máximo 1). Anexos rejeitados (nome, formato, conteúdo) ou
        com falha no download não geram log individual - só o resultado
        final da requisição importa para quem acompanha a extração (ver
        _processar).
        """
        total = len(elementos)
        salvos = 0
        arquivos_antes = len(self.arquivos_salvos_na_execucao)
        for i, el in enumerate(elementos, 1):
            if self.cancelado:
                break
            nome = await el.inner_text()

            if self.contem_de_acordo(nome):
                progress_down_callback(int((i / total) * 100))
                continue

            arquivo_temporario = os.path.join(self.pasta_download, f"temp_{nome}")
            try:
                async with page.expect_download(timeout=30000) as download_info:
                    await el.click()
                download = await download_info.value
                await download.save_as(arquivo_temporario)

                extensoes_suportadas = (".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".txt")
                if nome.lower().endswith(extensoes_suportadas):
                    sucesso, _motivo = self.analisar_arquivo(arquivo_temporario, req, nome)
                    if sucesso:
                        salvos += 1
            except Exception:
                pass
            finally:
                if os.path.exists(arquivo_temporario):
                    os.remove(arquivo_temporario)

            progress_down_callback(int((i / total) * 100))

        if salvos > 1:
            self._manter_apenas_melhor_candidato(self.arquivos_salvos_na_execucao[arquivos_antes:], req)
            salvos = len(self.arquivos_salvos_na_execucao) - arquivos_antes
        return salvos

    @staticmethod
    async def _particionar_anexos_carrinho(page, elementos):
        """Separa `elementos` entre os que pertencem a um item do carrinho
        (têm um ancestral .lineAttachments) e os da seção de Anexos de cima."""
        carrinho, topo = [], []
        for el in elementos:
            dentro_carrinho = await page.evaluate("(e) => e.closest('.lineAttachments') !== null", el)
            (carrinho if dentro_carrinho else topo).append(el)
        return carrinho, topo

    async def _processar(self, context, log_callback, progress_req_callback, progress_down_callback) -> bool:
        coupa_base_url = get_coupa_base_url()
        pages = context.pages
        page = pages[0] if pages else await context.new_page()

        total_reqs = len(self.requisicoes)
        for idx, req in enumerate(self.requisicoes, 1):
            if self.cancelado:
                break

            url = f"{coupa_base_url.rstrip('/')}/requisition_headers/{req.strip()}"
            log_callback(f"📂 Processando requisição #{req}...")

            try:
                try:
                    await page.goto(url, wait_until="load", timeout=60000)
                except Exception:
                    log_callback("⏳ Página demorou para carregar (possível tela de login)...")

                if "login" in page.url.lower() or "sso" in page.url.lower():
                    log_callback("⚠️ Realize o login no Edge se necessário...")
                    await page.wait_for_url(
                        lambda u: "login" not in u.lower() and "sso" not in u.lower(),
                        timeout=300000,
                    )
                    await page.wait_for_load_state("networkidle", timeout=15000)
                    await page.goto(url, wait_until="load", timeout=60000)

                aba_carrinho = await page.query_selector("a:has-text('Itens do carrinho')")
                if aba_carrinho:
                    await aba_carrinho.click()

                with contextlib.suppress(Exception):
                    # timeout aqui só significa "sem anexos" - a contagem logo
                    # abaixo (0 e 0) já deixa isso claro, sem precisar de um
                    # log extra com o texto cru da exceção.
                    await page.wait_for_selector(_SELETOR_ANEXOS, timeout=10000)

                todos_anexos = await page.query_selector_all(_SELETOR_ANEXOS)
                anexos_carrinho, anexos_topo = await self._particionar_anexos_carrinho(page, todos_anexos)
                log_callback(
                    f"🔎 #{req}: {len(anexos_carrinho)} anexo(s) no item do carrinho, "
                    f"{len(anexos_topo)} na seção de Anexos."
                )

                # 1. Anexo do item do carrinho primeiro - é sempre a versão mais
                # atualizada do orçamento. Só cai para a seção de Anexos (abaixo)
                # se o carrinho não tiver arquivo ou o arquivo não validar como orçamento.
                arquivos_salvos_no_req = 0
                if anexos_carrinho:
                    arquivos_salvos_no_req = await self._baixar_e_validar_anexos(
                        page, anexos_carrinho, req, progress_down_callback,
                    )

                if arquivos_salvos_no_req == 0 and anexos_topo:
                    arquivos_salvos_no_req = await self._baixar_e_validar_anexos(
                        page, anexos_topo, req, progress_down_callback,
                    )

                if arquivos_salvos_no_req > 0:
                    log_callback(f"✅ Orçamento salvo para a requisição #{req}.")
                else:
                    log_callback(f"❌ Nenhum orçamento válido encontrado para a requisição #{req}.")
                    self.requisicoes_sem_arquivos.append(req)

            except Exception as e:
                log_callback(f"❌ Erro ao processar #{req}: {str(e)}")
                log_callback(f"   Detalhes: {traceback.format_exc()}")

            progress_req_callback(int((idx / total_reqs) * 100))

        return not self.cancelado


class DownloadWorker(QThread):
    log_signal = pyqtSignal(str)
    progress_req_signal = pyqtSignal(int)
    progress_down_signal = pyqtSignal(int)
    finished_signal = pyqtSignal(bool, list, list)

    def __init__(self, requisicoes: list[str], pasta_download: str):
        super().__init__()
        self.scraper = DownloadScraper(requisicoes, pasta_download)
        # Melhoria 8: expõe o log_callback ao scraper para que extrair_texto
        # possa reportar erros via UI em vez de print().
        self.scraper._log_callback = self.log_signal.emit

    def run(self):
        """Melhoria 6: try/except garante que finished_signal SEMPRE seja emitido."""
        import asyncio
        import traceback as tb

        loop = None
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            sucesso = loop.run_until_complete(
                self.scraper.run(
                    self.log_signal.emit,
                    self.progress_req_signal.emit,
                    self.progress_down_signal.emit,
                )
            )
            self.finished_signal.emit(
                sucesso,
                self.scraper.arquivos_salvos_na_execucao,
                self.scraper.requisicoes_sem_arquivos,
            )
        except Exception as e:
            self.log_signal.emit(f"❌ ERRO CRÍTICO na thread de download: {str(e)}")
            with contextlib.suppress(Exception):
                # best-effort: não deixa o log do traceback mascarar o erro original
                self.log_signal.emit(tb.format_exc())
            # Garante que a UI não fique presa com botões desabilitados
            self.finished_signal.emit(False, [], self.scraper.requisicoes_sem_arquivos)
        finally:
            if loop is not None:
                with contextlib.suppress(Exception):
                    loop.close()  # best-effort: loop pode já estar fechado

    def cancelar(self) -> None:
        self.scraper.cancelado = True

